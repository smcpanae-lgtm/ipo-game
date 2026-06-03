"""
THE IPO PATH: 栄光への決断 — 概要説明書 PPT 生成スクリプト
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── カラーパレット（ゲームのUIカラーに合わせた紺系） ──
C_NAVY      = RGBColor(0x0A, 0x1D, 0x38)   # 濃紺（背景）
C_MID_NAVY  = RGBColor(0x12, 0x2B, 0x52)   # 中紺
C_BLUE      = RGBColor(0x1A, 0x56, 0x8C)   # メインブルー
C_ACCENT    = RGBColor(0x3A, 0xA0, 0xE8)   # アクセントブルー
C_GOLD      = RGBColor(0xE8, 0xC0, 0x00) # ゴールド（近似）
C_GOLD2     = RGBColor(0xE8, 0xB8, 0x40)   # ゴールド
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xD0, 0xE8, 0xFF)   # 薄青
C_RED       = RGBColor(0xE8, 0x40, 0x40)
C_GREEN     = RGBColor(0x40, 0xC0, 0x80)
C_ORANGE    = RGBColor(0xF0, 0x90, 0x20)
C_GRAY_LT   = RGBColor(0xB0, 0xC8, 0xE0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # 完全ブランク
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             font_name="Meiryo UI", italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  font_size=14, color=None, bold_first=False,
                  line_spacing=1.2, font_name="Meiryo UI"):
    """lines = list of (text, bold, color_override)  or str"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bld, col = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = _Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.size = _Pt(font_size)
        run.font.bold = bld
        run.font.name = font_name
        if col:
            run.font.color.rgb = col
    return txBox


def header_bar(slide, title_text, accent_color=None):
    """スライドヘッダーバー共通部品"""
    ac = accent_color or C_ACCENT
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=C_NAVY)
    add_rect(slide, 0, 1.0, 13.33, 0.06, fill_color=ac)
    add_text(slide, "THE IPO PATH : 栄光への決断",
             0.25, 0.04, 6, 0.45,
             font_size=11, color=C_GRAY_LT, bold=False)
    add_text(slide, title_text,
             0.25, 0.42, 12.5, 0.65,
             font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def footer_bar(slide, page_num, total=10):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=C_NAVY)
    add_text(slide, "© THE IPO PATH: 栄光への決断  —  IPO検定協会 ご説明資料",
             0.3, 7.12, 10, 0.36, font_size=10, color=C_GRAY_LT)
    add_text(slide, f"{page_num} / {total}",
             12.5, 7.12, 0.7, 0.36, font_size=10, color=C_GRAY_LT,
             align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════
#  SLIDE 1 : 表紙
# ════════════════════════════════════════════════════════
def slide01_cover(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, C_NAVY)

    # グラデーション的な装飾矩形
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_MID_NAVY)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD2)
    add_rect(slide, 0, 7.42, 13.33, 0.08, fill_color=C_GOLD2)
    # 左縦ストライプ装飾
    add_rect(slide, 0, 0, 0.18, 7.5, fill_color=C_BLUE)
    add_rect(slide, 0.18, 0, 0.05, 7.5, fill_color=C_ACCENT)

    # サブタイトル帯
    add_rect(slide, 1.0, 2.65, 11.1, 0.07, fill_color=C_ACCENT)

    # メインタイトル
    add_text(slide, "THE IPO PATH",
             1.2, 1.2, 11, 1.05,
             font_size=54, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "栄 光 へ の 決 断",
             1.2, 2.1, 11, 0.65,
             font_size=28, bold=True, color=C_GOLD2, align=PP_ALIGN.LEFT)

    # 説明資料ラベル
    add_text(slide, "ゲーム概要説明書・仕様書",
             1.2, 2.85, 11, 0.55,
             font_size=22, bold=False, color=C_LIGHT, align=PP_ALIGN.LEFT)

    # 概要説明
    add_multiline(slide, [
        ("IPO（株式公開）準備のリアルな意思決定を体験する教育型シミュレーションゲーム", False, C_LIGHT),
        ("日本公認会計士協会「新規上場のための事前準備ガイドブック」準拠", False, C_GRAY_LT),
    ], 1.2, 3.55, 11, 0.8, font_size=15)

    # バッジ
    add_rect(slide, 1.2, 4.55, 4.5, 0.55, fill_color=C_BLUE, line_color=C_ACCENT, line_width_pt=1.5)
    add_text(slide, "👤  プレイヤーの役割：代表取締役社長",
             1.3, 4.58, 4.3, 0.48, font_size=13, bold=True, color=C_WHITE)

    add_rect(slide, 6.0, 4.55, 3.5, 0.55, fill_color=C_BLUE, line_color=C_ACCENT, line_width_pt=1.5)
    add_text(slide, "🌐  形式：Webブラウザ型ゲーム",
             6.1, 4.58, 3.3, 0.48, font_size=13, bold=True, color=C_WHITE)

    add_rect(slide, 9.8, 4.55, 3.3, 0.55, fill_color=C_BLUE, line_color=C_ACCENT, line_width_pt=1.5)
    add_text(slide, "📚  対象：IPO実務者・学習者",
             9.9, 4.58, 3.1, 0.48, font_size=13, bold=True, color=C_WHITE)

    # 提供先
    add_text(slide, "提供先：IPO検定協会  様",
             1.2, 5.5, 11, 0.5, font_size=16, bold=True, color=C_GOLD2)
    add_text(slide, "© THE IPO PATH: 栄光への決断",
             1.2, 6.8, 11, 0.45, font_size=11, color=C_GRAY_LT)


# ════════════════════════════════════════════════════════
#  SLIDE 2 : ゲームコンセプト・目的
# ════════════════════════════════════════════════════════
def slide02_concept(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "ゲームのコンセプト・目的")
    footer_bar(slide, 2)

    # 左カラム：目的
    add_rect(slide, 0.3, 1.3, 5.9, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_ACCENT, line_width_pt=1.2)
    add_rect(slide, 0.3, 1.3, 5.9, 0.42, fill_color=C_BLUE)
    add_text(slide, "🎯  開発の目的", 0.45, 1.32, 5.6, 0.38,
             font_size=14, bold=True, color=C_WHITE)
    add_multiline(slide, [
        ("IPO準備の実務知識を「体験しながら自然に習得」する", True, C_LIGHT),
        "",
        ("● 公認会計士協会ガイドブックの実務指摘事項を", False, C_LIGHT),
        ("  ゲームギミックとして落とし込み", False, C_LIGHT),
        ("● 選択の結果が将来の上場審査に影響する", False, C_LIGHT),
        ("  「因果応報」システムで実務感覚を養成", False, C_LIGHT),
        ("● 40を超える意思決定場面を通じて", False, C_LIGHT),
        ("  IPO実務の全体像を体系的に学べる", False, C_LIGHT),
        ("● Google Gemini AI がナラティブを動的生成", False, RGBColor(0x80, 0xD0, 0xFF)),
        ("  毎ターン・各意思決定に応じた物語をAIが生成", False, RGBColor(0x80, 0xD0, 0xFF)),
        "",
        ("■ 準拠資料", True, C_GOLD2),
        ("日本公認会計士協会", False, C_LIGHT),
        ("「新規上場のための事前準備ガイドブック」", True, C_LIGHT),
    ], 0.45, 1.82, 5.6, 4.8, font_size=13)

    # 右カラム：ターゲット
    add_rect(slide, 6.5, 1.3, 6.5, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_ACCENT, line_width_pt=1.2)
    add_rect(slide, 6.5, 1.3, 6.5, 0.42, fill_color=C_BLUE)
    add_text(slide, "👥  想定利用者・活用場面", 6.65, 1.32, 6.2, 0.38,
             font_size=14, bold=True, color=C_WHITE)
    targets = [
        ("IPO準備中のスタートアップ経営者・CFO", C_LIGHT),
        ("証券会社・監査法人のIPO担当者", C_LIGHT),
        ("IPO検定の受験者・学習者", C_GOLD2),
        ("経営管理部門・管理部門スタッフ", C_LIGHT),
        ("ベンチャーキャピタル・投資家", C_LIGHT),
    ]
    for i, (txt, col) in enumerate(targets):
        y = 1.9 + i * 0.62
        add_rect(slide, 6.65, y, 6.2, 0.5,
                 fill_color=RGBColor(0x16, 0x34, 0x60),
                 line_color=RGBColor(0x30, 0x60, 0x90), line_width_pt=0.8)
        add_text(slide, f"  ✦  {txt}", 6.7, y+0.06, 6.1, 0.38,
                 font_size=13, color=col)

    add_rect(slide, 6.65, 5.15, 6.2, 1.5, fill_color=RGBColor(0x16, 0x34, 0x60),
             line_color=C_GOLD2, line_width_pt=1.2)
    add_text(slide, "💡  活用場面", 6.8, 5.18, 5.8, 0.38,
             font_size=13, bold=True, color=C_GOLD2)
    add_multiline(slide, [
        "研修・勉強会のインタラクティブ教材",
        "IPO検定学習補助ツール",
        "社内IPO勉強会のコンテンツ",
    ], 6.8, 5.55, 6.0, 1.0, font_size=12, color=C_LIGHT)


# ════════════════════════════════════════════════════════
#  SLIDE 3 : ゲーム全体構造
# ════════════════════════════════════════════════════════
def slide03_structure(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "ゲームの全体構造・タイムライン")
    footer_bar(slide, 3)

    # タイムライン矢印帯
    add_rect(slide, 0.3, 1.35, 12.73, 0.65, fill_color=C_NAVY)

    periods = [
        ("N-3期\n（体制構築期）", C_BLUE,       0.35),
        ("N-2期\n（直前々期）",   C_MID_NAVY,   3.55),
        ("N-1期\n（直前期）",     C_BLUE,       6.75),
        ("N期\n（申請・上場）",   RGBColor(0x8B, 0x20, 0x20), 9.95),
    ]
    for label, col, x in periods:
        add_rect(slide, x, 1.32, 3.05, 0.7, fill_color=col,
                 line_color=C_ACCENT, line_width_pt=0.8)
        add_text(slide, label, x+0.05, 1.34, 2.9, 0.65,
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 矢印
    for ax in [3.38, 6.58, 9.78]:
        add_text(slide, "▶", ax, 1.44, 0.22, 0.45,
                 font_size=18, bold=True, color=C_ACCENT)

    # 各期の詳細ボックス
    details = [
        ("N-3期", [
            "● ショートレビュー実施",
            "● 監査法人候補の選定",
            "● CFO・管理体制構築",
            "● 主幹事証券会社の選定",
            "● 中期経営計画策定",
            "● 発生主義・証憑管理整備",
        ], 0.35),
        ("N-2期", [
            "● 監査契約締結（2期監査開始）",
            "● 内部管理体制の運用開始",
            "● 職務権限規程・予算管理",
            "● 棚卸立会の実施",
            "● 関連当事者取引の整理",
            "● シリーズBファイナンス",
        ], 3.55),
        ("N-1期", [
            "● 内部統制報告制度（J-SOX）",
            "● 会計監査人の正式選任",
            "● インサイダー取引防止",
            "● 適時開示体制の整備",
            "● 機関投資家ロードショー",
            "● 有価証券届出書の準備",
        ], 6.75),
        ("N期", [
            "● 上場申請書類の提出",
            "● 東証実質審査・ヒアリング",
            "● 公募・売出し規模の決定",
            "● 公開価格・需要調査",
            "● 上場承認・株式公開",
            "● 東証の鐘を鳴らせ！",
        ], 9.95),
    ]
    for title, items, x in details:
        add_rect(slide, x, 2.15, 3.05, 4.7,
                 fill_color=RGBColor(0x10, 0x28, 0x48),
                 line_color=C_ACCENT, line_width_pt=0.8)
        for j, item in enumerate(items):
            col = C_GOLD2 if j == 5 and title == "N期" else C_LIGHT
            add_text(slide, item, x+0.12, 2.2 + j*0.72, 2.85, 0.68,
                     font_size=11.5, color=col)

    # 下部補足
    add_rect(slide, 0.3, 6.98, 12.73, 0.45, fill_color=C_NAVY)
    add_text(slide,
             "1ターン = 1四半期（Q）　／　全16ターン（N-3〜N 各4Q）　／　各Qに1〜2個のイベント発生",
             0.5, 6.99, 12.5, 0.42,
             font_size=12.5, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 4 : ゲームシステム（スコア・フラグ）
# ════════════════════════════════════════════════════════
def slide04_system(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "ゲームシステム — スコア・フラグ・因果応報")
    footer_bar(slide, 4)

    # 左：7つのスコア
    add_rect(slide, 0.3, 1.3, 5.9, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_ACCENT, line_width_pt=1.2)
    add_rect(slide, 0.3, 1.3, 5.9, 0.42, fill_color=C_BLUE)
    add_text(slide, "📊  管理スコア 7種（各0〜100）", 0.45, 1.32, 5.6, 0.38,
             font_size=14, bold=True, color=C_WHITE)

    scores = [
        ("内部管理体制スコア",  "各種規程・職務分掌の整備状況",     C_ACCENT),
        ("コンプライアンス",    "法令遵守・反社対応・労務管理",      C_ACCENT),
        ("決算品質スコア",      "会計基準・証憑管理・棚卸立会",      C_ACCENT),
        ("ガバナンススコア",    "社外役員・取締役会の実効性",        C_ACCENT),
        ("監査法人信頼度",      "監査法人との信頼関係",              C_GOLD2),
        ("投資家信頼度",        "証券会社・投資家からの評価",        C_GOLD2),
        ("従業員士気",          "組織・人材マネジメント",            C_GREEN),
    ]
    for i, (name, desc, col) in enumerate(scores):
        y = 1.85 + i * 0.67
        add_rect(slide, 0.42, y, 5.65, 0.58,
                 fill_color=RGBColor(0x14, 0x30, 0x58),
                 line_color=RGBColor(0x30, 0x60, 0x90), line_width_pt=0.6)
        add_text(slide, name, 0.55, y+0.04, 2.5, 0.28,
                 font_size=12, bold=True, color=col)
        add_text(slide, desc, 0.55, y+0.28, 5.3, 0.25,
                 font_size=10.5, color=C_GRAY_LT)

    # 右：フラグ・因果応報
    add_rect(slide, 6.5, 1.3, 6.5, 2.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=RGBColor(0xC0, 0x30, 0x30), line_width_pt=1.2)
    add_rect(slide, 6.5, 1.3, 6.5, 0.42, fill_color=RGBColor(0x8B, 0x20, 0x20))
    add_text(slide, "💣  因果応報フラグ（隠れ爆弾）", 6.65, 1.32, 6.2, 0.38,
             font_size=14, bold=True, color=C_WHITE)
    bombs = [
        "⚠️  未払残業代 → 労基署調査（カウントダウン）",
        "💣  反社チェック不備 → 主幹事からの通告",
        "🔓  職務未分掌 → 横領リスクの蓄積",
        "📉  現金主義会計 → 発生主義移行未完了",
        "🔥  利益操作 → 監査発覚タイマー",
    ]
    for i, b in enumerate(bombs):
        add_text(slide, b, 6.65, 1.82 + i*0.39, 6.2, 0.35,
                 font_size=11, color=RGBColor(0xFF, 0xB0, 0xB0))

    add_rect(slide, 6.5, 3.97, 6.5, 2.88, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_GOLD2, line_width_pt=1.2)
    add_rect(slide, 6.5, 3.97, 6.5, 0.42, fill_color=RGBColor(0x5B, 0x4A, 0x10))
    add_text(slide, "⚖️  上場審査 リスクスコア", 6.65, 3.99, 6.2, 0.38,
             font_size=14, bold=True, color=C_GOLD2)
    add_multiline(slide, [
        ("60以上 → 審査チェック1項目NG（他と重なると不通過）", True, C_GOLD2),
        ("80以上（N期中）→ 即・上場計画全面撤回エンディング", True, RGBColor(0xFF, 0x60, 0x60)),
        "N-3期の手抜きがN期の否認を招く「因果応報」設計",
        "ショートレビューで隠れリスクが「可視化」される",
        "各判断の影響が複数期間にわたって連鎖する",
    ], 6.65, 4.45, 6.2, 2.3, font_size=11.5, color=C_LIGHT)


# ════════════════════════════════════════════════════════
#  SLIDE 5 : イベントシステム（56種）
# ════════════════════════════════════════════════════════
def slide05_events(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "イベントシステム — 56種類の意思決定場面")
    footer_bar(slide, 5)

    # 上部：イベント概要
    add_rect(slide, 0.3, 1.3, 12.73, 0.85, fill_color=C_NAVY)
    add_text(slide,
             "各ターンに1〜2個のIPO実務イベントが発生。プレイヤーが選択肢A〜Cを選び、その影響が即時・将来の両方に波及する。",
             0.5, 1.35, 12.4, 0.75, font_size=13, color=C_LIGHT)

    # カテゴリ別イベント表
    categories = [
        ("N-3期 主要イベント（体制構築）", C_BLUE, [
            "ショートレビューの実施",
            "監査法人候補の選定",
            "CFO採用",
            "主幹事証券会社の選定",
            "中期経営計画策定",
            "発生主義・収益認識移行",
            "証憑管理の整備",
            "棚卸資産・原価計算体制",
        ]),
        ("N-2期 主要イベント（整備・運用）", C_BLUE, [
            "監査契約締結・ショートレビュー対応",
            "シリーズB大型ファイナンス",
            "月次決算早期化（10日締め）",
            "職務権限・業務分掌規程",
            "予算管理・内部監査設置",
            "関連当事者取引の整理",
            "独立社外役員の選任",
            "J-SOX内部統制整備",
        ]),
        ("N-1〜N期 主要イベント（審査・上場）", C_BLUE, [
            "会計監査人の正式選任（AGM）",
            "インサイダー取引防止体制",
            "適時開示・IR体制整備",
            "機関投資家ロードショー",
            "有価証券届出書の準備",
            "SO行使・公募売出し規模",
            "公開価格・需要調査",
            "上場最終手続き（ほふり等）",
        ]),
    ]
    col_x = [0.3, 4.6, 8.9]
    for ci, (cat_title, col, items) in enumerate(categories):
        x = col_x[ci]
        add_rect(slide, x, 2.25, 4.15, 4.65,
                 fill_color=RGBColor(0x10, 0x28, 0x48),
                 line_color=C_ACCENT, line_width_pt=0.8)
        add_rect(slide, x, 2.25, 4.15, 0.38, fill_color=C_BLUE)
        add_text(slide, cat_title, x+0.1, 2.27, 3.95, 0.35,
                 font_size=11.5, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            add_text(slide, f"  • {item}", x+0.1, 2.72 + j*0.51, 3.95, 0.48,
                     font_size=11, color=C_LIGHT)

    # 下部：イベント数バッジ
    add_rect(slide, 0.3, 7.02, 12.73, 0.38, fill_color=C_NAVY)
    add_text(slide,
             "全56イベント  |  ワールドイベント（地政学・規制変更・金利ショック等）も発生  |  毎プレイで順序が変化するリプレイ性",
             0.5, 7.03, 12.4, 0.35,
             font_size=11.5, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 6 : 学習コンテンツ（IPO実務10論点）
# ════════════════════════════════════════════════════════
def slide06_learning(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "学習コンテンツ — IPO実務10論点への対応")
    footer_bar(slide, 6)

    add_text(slide, "公認会計士協会ガイドブック「上場審査の10論点」をゲームに完全組み込み",
             0.5, 1.28, 12.3, 0.38, font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

    points = [
        ("① ガバナンス・コンプライアンス",   "社外役員、反社排除、内部通報制度"),
        ("② 会計・内部管理体制",              "発生主義、証憑管理、月次決算早期化"),
        ("③ 収益・売上の信頼性",              "収益認識基準、顧客集中リスク対応"),
        ("④ 原価計算・棚卸資産",              "棚卸立会、原価計算体制の構築"),
        ("⑤ 財務・資本政策",                  "キャップテーブル、種類株式転換、資本政策"),
        ("⑥ 関連当事者取引",                  "オーナー取引の整理・開示"),
        ("⑦ 組織・人材・キーパーソン",        "CFO採用、職務分掌、キーパーソン依存"),
        ("⑧ 知財・IT・情報管理",              "IPO保護体制、ERPシステム整備"),
        ("⑨ 監査・内部統制（J-SOX）",         "監査契約、3点セット整備、内部監査"),
        ("⑩ 証券市場・適時開示",              "IR体制、ロードショー、適時開示訓練"),
    ]
    for i, (title, detail) in enumerate(points):
        row = i // 2
        col = i % 2
        x = 0.3 + col * 6.5
        y = 1.8 + row * 0.98
        add_rect(slide, x, y, 6.2, 0.87,
                 fill_color=RGBColor(0x10, 0x28, 0x48),
                 line_color=C_ACCENT, line_width_pt=0.7)
        add_text(slide, title, x+0.12, y+0.04, 5.9, 0.35,
                 font_size=12.5, bold=True, color=C_GOLD2)
        add_text(slide, detail, x+0.12, y+0.38, 5.9, 0.4,
                 font_size=11.5, color=C_LIGHT)

    add_rect(slide, 0.3, 6.82, 12.73, 0.55, fill_color=C_NAVY)
    add_text(slide,
             "💡 各イベントにはIPO先生（AIアドバイザー）のワンポイント解説ボタンが付属 — 実務知識を即座に確認可能",
             0.5, 6.84, 12.4, 0.5, font_size=13, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 7 : 画面構成・UI
# ════════════════════════════════════════════════════════
def slide07_ui(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "画面構成・ユーザーインターフェース")
    footer_bar(slide, 7)

    # メイン画面モック（3カラム構造）
    add_rect(slide, 0.3, 1.3, 12.73, 5.6, fill_color=RGBColor(0x06, 0x10, 0x22),
             line_color=C_ACCENT, line_width_pt=1.5)

    # ヘッダー
    add_rect(slide, 0.3, 1.3, 12.73, 0.45, fill_color=RGBColor(0x08, 0x18, 0x35))
    add_text(slide, "🏛 THE IPO PATH  ⏰ N-3期 Q2  |  💾 セーブ  |  💡 IPO先生",
             0.5, 1.33, 12.3, 0.38, font_size=11, color=C_GOLD2)

    # 左サイドバー
    add_rect(slide, 0.3, 1.75, 2.45, 5.15, fill_color=RGBColor(0x0C, 0x1E, 0x3C),
             line_color=RGBColor(0x1A, 0x40, 0x70), line_width_pt=0.8)
    add_rect(slide, 0.3, 1.75, 2.45, 0.32, fill_color=RGBColor(0x1A, 0x40, 0x70))
    add_text(slide, "📊 会社ダッシュボード", 0.4, 1.76, 2.25, 0.3,
             font_size=9, bold=True, color=C_WHITE)
    left_items = [
        "💰 手元資金: ¥680M",
        "📈 売上: ¥120M/Q",
        "🔥 ランウェイ: 12Q",
        "─────────────────",
        "内部管理体制:  ████ 62",
        "コンプライアンス:████ 58",
        "決算品質:      ██ 45",
        "ガバナンス:    ████ 70",
        "─────────────────",
        "監査法人信頼:  ████ 72",
        "投資家信頼:    ███ 65",
        "従業員士気:    ████ 75",
        "─────────────────",
        "⚠️ リスクスコア: 28",
    ]
    for i, item in enumerate(left_items):
        col = C_LIGHT if not item.startswith("⚠️") else RGBColor(0xFF, 0xA0, 0xA0)
        add_text(slide, item, 0.4, 2.12 + i*0.3, 2.3, 0.28,
                 font_size=8.5, color=col)

    # メインストーリー
    add_rect(slide, 2.82, 1.75, 7.5, 5.15, fill_color=RGBColor(0x0A, 0x1C, 0x38))
    add_rect(slide, 2.82, 1.75, 7.5, 0.32, fill_color=RGBColor(0x12, 0x2B, 0x52))
    add_text(slide, "📖 メインストーリー（意思決定パネル）", 2.95, 1.77, 7.2, 0.28,
             font_size=9, bold=True, color=C_WHITE)
    # ストーリーテキスト
    add_rect(slide, 2.9, 2.13, 7.3, 1.1, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=RGBColor(0x30, 0x60, 0x90), line_width_pt=0.5)
    add_text(slide, "📋 主幹事証券会社の選定・引受契約",
             3.0, 2.15, 7.0, 0.35, font_size=10.5, bold=True, color=C_GOLD2)
    add_text(slide, "CFOからの報告：「主幹事証券会社の選定について、タイミングが極めて重要です。\n（N-3期が理想、N-2期が最遅）公開指導期間の確保が上場成功の鍵です…」",
             3.0, 2.52, 7.1, 0.65, font_size=9, color=C_LIGHT)
    # 選択肢
    choice_colors = [C_GREEN, C_BLUE, RGBColor(0x80, 0x50, 0x10)]
    choices = [
        "A. 大手主幹事証券会社を選定する（審査厳格・ネットワーク最強）",
        "B. 中堅証券会社を選定する（機動的・コスト効率重視）",
        "C. 今期は選定を見送り、N-1期以降に対応する  ⚠️ 遅延リスク",
    ]
    for i, (ch, col) in enumerate(zip(choices, choice_colors)):
        add_rect(slide, 2.9, 3.34 + i * 0.62, 7.3, 0.52,
                 fill_color=RGBColor(0x14, 0x30, 0x58),
                 line_color=col, line_width_pt=1.0)
        add_text(slide, ch, 3.0, 3.38 + i * 0.62, 7.1, 0.42,
                 font_size=10, color=C_LIGHT)
    add_text(slide, "※ ダブルクリックで選択確定 / 1回クリックで「もう一度クリックで決定」",
             2.9, 5.28, 7.3, 0.45, font_size=9, color=C_GRAY_LT, italic=True)

    # 右サイドバー
    add_rect(slide, 10.39, 1.75, 2.64, 5.15, fill_color=RGBColor(0x0C, 0x1E, 0x3C),
             line_color=RGBColor(0x1A, 0x40, 0x70), line_width_pt=0.8)
    add_rect(slide, 10.39, 1.75, 2.64, 0.32, fill_color=RGBColor(0x1A, 0x40, 0x70))
    add_text(slide, "✅ 上場準備チェックリスト", 10.48, 1.77, 2.45, 0.28,
             font_size=9, bold=True, color=C_WHITE)
    checks = [
        ("✅", "ショートレビュー", C_GREEN),
        ("✅", "監査法人候補選定", C_GREEN),
        ("✅", "CFO採用", C_GREEN),
        ("⬜", "主幹事証券会社選定", C_ORANGE),
        ("⬜", "中期経営計画", C_GRAY_LT),
        ("⬜", "社外役員選任", C_GRAY_LT),
        ("⬜", "監査契約締結", C_GRAY_LT),
        ("⬜", "発生主義移行", C_GRAY_LT),
    ]
    for i, (mark, txt, col) in enumerate(checks):
        add_text(slide, f"{mark} {txt}", 10.5, 2.14 + i * 0.53, 2.45, 0.48,
                 font_size=9.5, color=col)

    # ラベル注釈
    labels = [
        (0.3,  6.93, "① 左サイドバー\n会社状況ダッシュボード"),
        (3.8,  6.93, "② メインストーリー\n意思決定・選択肢パネル"),
        (10.5, 6.93, "③ 右サイドバー\n上場準備チェックリスト"),
    ]
    for lx, ly, ltxt in labels:
        add_text(slide, ltxt, lx, ly, 3.0, 0.5,
                 font_size=10, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 8 : 操作方法
# ════════════════════════════════════════════════════════
def slide08_howto(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "操作方法・ゲームの進め方")
    footer_bar(slide, 8)

    steps = [
        ("STEP 1", "ゲーム開始",
         "タイトル画面「▶ 栄光への旅を始める」をクリック\n業種（SaaS / FinTech / 製造業 / 小売業）を選択\n市場区分（グロース / スタンダード / プライム）を選択\n会社名・代表者名を入力してゲームスタート",
         C_ACCENT),
        ("STEP 2", "ターン進行",
         "各ターン（四半期）にIPO実務イベントが1〜2件表示\n「IPO先生」ボタンでイベントの実務背景を確認可能\n選択肢A〜C（またはA〜D）から判断を選ぶ\n「ダブルクリック」で選択確定（1クリックは確認段階）",
         C_ACCENT),
        ("STEP 3", "結果確認",
         "選択後にルーレット演出と結果が表示される\n各スコアの変化・フラグの変動を確認\n「続ける▶」ボタンで次のターンへ進む\nセーブボタン（💾）でいつでもデータを保存可能",
         C_ACCENT),
        ("STEP 4", "上場審査（クライマックス）",
         "N期に東証への上場申請・実質審査が始まる\nリスクスコア100未満・各種要件充足が審査通過条件\nN-3期からの判断の積み重ねが審査結果を左右する\n「東証の鐘」を鳴らしてエンディングへ！",
         C_GOLD2),
    ]
    for i, (step, title, desc, col) in enumerate(steps):
        row = i // 2
        c_ = i % 2
        x = 0.3 + c_ * 6.5
        y = 1.28 + row * 2.85
        add_rect(slide, x, y, 6.2, 2.65,
                 fill_color=RGBColor(0x10, 0x28, 0x48),
                 line_color=col, line_width_pt=1.5)
        add_rect(slide, x, y, 1.1, 2.65, fill_color=C_NAVY)
        add_text(slide, step, x+0.05, y+0.7, 1.0, 0.9,
                 font_size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+1.2, y+0.06, 4.85, 0.45,
                 font_size=15, bold=True, color=col)
        add_text(slide, desc, x+1.2, y+0.55, 4.85, 1.9,
                 font_size=11, color=C_LIGHT)

    # ショートカット
    add_rect(slide, 0.3, 7.0, 12.73, 0.38, fill_color=C_NAVY)
    add_text(slide,
             "💡  再開：タイトル画面「続きから再開」  ／  保存：ヘッダーの💾ボタン  ／  BGM：画面右下SOUNDボタン",
             0.5, 7.01, 12.4, 0.35, font_size=11.5, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  SLIDE 9 : 特徴・差別化ポイント
# ════════════════════════════════════════════════════════
def slide09_features(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "ゲームの特長・差別化ポイント")
    footer_bar(slide, 9)

    features = [
        ("🎓", "実務準拠の学習設計",
         "公認会計士協会「事前準備ガイドブック」の\n指摘事項をゲームに完全組み込み\n全56イベントが実際のIPO実務を反映",
         C_ACCENT),
        ("💣", "因果応報システム",
         "N-3期の手抜きがN期の審査否認を招く\n「隠れ爆弾」フラグが将来の上場を直撃\n選択の影響が複数期間にわたって連鎖",
         C_RED),
        ("🎯", "リアルなIPOスケジュール",
         "N-3〜N期の4年間を16ターンで再現\n定時株主総会（AGM）システムを搭載\n会社法・金商法の手続きタイミングも反映",
         C_GREEN),
        ("📊", "多彩なゲームパラメータ",
         "7種類のスコア + リスクスコアで会社状態を管理\n業種4種 × 市場区分3種 = 12通りの組み合わせ\nワールドイベントで外部環境変化も体験",
         C_GOLD2),
        ("🔄", "高いリプレイ性",
         "毎プレイでイベント発生順が変化\n複数の選択肢が異なる将来結果を生む\nセーブ＆ロードで選択肢の比較学習が可能",
         C_ACCENT),
        ("🌐", "ブラウザ型・導入不要",
         "Webブラウザで即座にプレイ開始\nインストール・ダウンロード不要\nスマートフォン・タブレットにも対応",
         C_GREEN),
    ]
    # 上段6件：2行 × 3列（やや縦短め）
    for i, (icon, title, desc, col) in enumerate(features):
        row = i // 3
        c_ = i % 3
        x = 0.3 + c_ * 4.35
        y = 1.3 + row * 2.5
        add_rect(slide, x, y, 4.1, 2.3,
                 fill_color=RGBColor(0x10, 0x28, 0x48),
                 line_color=col, line_width_pt=1.5)
        add_text(slide, icon, x+0.15, y+0.1, 0.7, 0.55,
                 font_size=20, align=PP_ALIGN.CENTER)
        add_text(slide, title, x+0.85, y+0.1, 3.1, 0.42,
                 font_size=13, bold=True, color=col)
        add_text(slide, desc, x+0.15, y+0.58, 3.8, 1.6,
                 font_size=10.5, color=C_LIGHT)

    # 7番目：Google Gemini AI — 全幅バナーカード
    ai_col = RGBColor(0x40, 0xB0, 0xFF)
    add_rect(slide, 0.3, 6.42, 12.73, 0.95,
             fill_color=RGBColor(0x08, 0x20, 0x44),
             line_color=ai_col, line_width_pt=2.0)
    add_text(slide, "🤖", 0.5, 6.47, 0.7, 0.75,
             font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide, "Google Gemini AI 搭載",
             1.25, 6.46, 3.5, 0.38,
             font_size=14, bold=True, color=ai_col)
    add_text(slide,
             "毎ターン冒頭の状況描写・突発クライシスのシナリオ・意思決定結果のドラマ演出を"
             "Google Gemini AIがリアルタイムで動的生成。毎プレイで異なるナラティブ体験を実現。"
             "（使用モデル：gemini-2.5-flash-lite）",
             1.25, 6.84, 11.6, 0.48,
             font_size=11, color=C_LIGHT)


# ════════════════════════════════════════════════════════
#  SLIDE 10 : 技術仕様・今後の展開
# ════════════════════════════════════════════════════════
def slide10_tech(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, RGBColor(0x0E, 0x20, 0x3A))
    header_bar(slide, "技術仕様・動作環境・今後の展開")
    footer_bar(slide, 10)

    # 左：技術仕様
    add_rect(slide, 0.3, 1.3, 5.9, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_ACCENT, line_width_pt=1.2)
    add_rect(slide, 0.3, 1.3, 5.9, 0.42, fill_color=C_BLUE)
    add_text(slide, "🛠  技術仕様", 0.45, 1.32, 5.6, 0.38,
             font_size=14, bold=True, color=C_WHITE)
    specs = [
        ("言語・フレームワーク", "Python 3.10+ / Flask"),
        ("フロントエンド",       "HTML5 / CSS3 / Vanilla JS"),
        ("ゲームエンジン",       "独自設計（Pythonクラスベース）"),
        ("データモデル",         "Python dataclasses"),
        ("AIエンジン",           "Google Gemini API"),
        ("AIモデル",             "gemini-2.5-flash-lite"),
        ("セーブデータ",         "dill + localStorage"),
        ("BGM",                  "Web Audio API（コード生成）"),
        ("動作確認ブラウザ",     "Chrome / Edge / Firefox"),
        ("推奨画面解像度",       "1280×720以上"),
    ]
    for i, (k, v) in enumerate(specs):
        y = 1.85 + i * 0.47
        add_text(slide, k, 0.45, y, 2.2, 0.42, font_size=11, color=C_GRAY_LT)
        add_text(slide, v, 2.7, y, 3.3, 0.42, font_size=11, bold=True, color=C_LIGHT)

    # 中：現状スペック
    add_rect(slide, 6.5, 1.3, 3.0, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_GOLD2, line_width_pt=1.2)
    add_rect(slide, 6.5, 1.3, 3.0, 0.42, fill_color=RGBColor(0x5B, 0x4A, 0x10))
    add_text(slide, "📋  現状スペック", 6.65, 1.32, 2.7, 0.38,
             font_size=14, bold=True, color=C_GOLD2)
    current_specs = [
        ("全イベント数",     "56種"),
        ("ターン数",         "最大16Q（4年）"),
        ("業種",             "4種"),
        ("市場区分",         "3種"),
        ("管理スコア",       "7種"),
        ("エンディング",     "複数パターン"),
        ("ゲームクリア条件", "東証上場承認"),
        ("ゲームオーバー",   "資金枯渇\n上場審査否認\nN期終了"),
    ]
    for i, (k, v) in enumerate(current_specs):
        y = 1.85 + i * 0.57
        add_text(slide, k, 6.65, y, 1.5, 0.5, font_size=11, color=C_GRAY_LT)
        add_text(slide, v, 8.2, y, 1.2, 0.5, font_size=11, bold=True, color=C_GOLD2)

    # 右：今後の展開
    add_rect(slide, 9.7, 1.3, 3.63, 5.55, fill_color=RGBColor(0x10, 0x28, 0x48),
             line_color=C_GREEN, line_width_pt=1.2)
    add_rect(slide, 9.7, 1.3, 3.63, 0.42, fill_color=RGBColor(0x10, 0x4A, 0x28))
    add_text(slide, "🚀  今後の展開案", 9.85, 1.32, 3.4, 0.38,
             font_size=14, bold=True, color=C_GREEN)
    roadmap = [
        ("クラウド展開",         "Webサーバーへの\nデプロイで複数人同時利用"),
        ("マルチプレイ",         "グループ研修対応\nチーム意思決定モード"),
        ("スコア集計",           "研修参加者の\nスコア・選択履歴を集計"),
        ("イベント追加",         "IPO検定出題範囲に\n合わせた問題の拡充"),
        ("解説資料連携",         "各イベントと\n教材資料のリンク機能"),
    ]
    for i, (k, v) in enumerate(roadmap):
        y = 1.85 + i * 0.95
        add_rect(slide, 9.85, y, 3.35, 0.8,
                 fill_color=RGBColor(0x14, 0x38, 0x24),
                 line_color=RGBColor(0x20, 0x80, 0x48), line_width_pt=0.7)
        add_text(slide, k, 9.98, y+0.04, 1.3, 0.32,
                 font_size=11, bold=True, color=C_GREEN)
        add_text(slide, v, 9.98, y+0.38, 3.1, 0.38,
                 font_size=10, color=C_LIGHT)

    # フッター強調
    add_rect(slide, 0.3, 6.98, 12.73, 0.4, fill_color=C_NAVY)
    add_text(slide,
             "ご質問・カスタマイズのご要望はお気軽にお申し付けください",
             0.5, 6.99, 12.4, 0.38,
             font_size=14, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide01_cover(prs)
    slide02_concept(prs)
    slide03_structure(prs)
    slide04_system(prs)
    slide05_events(prs)
    slide06_learning(prs)
    slide07_ui(prs)
    slide08_howto(prs)
    slide09_features(prs)
    slide10_tech(prs)

    import io
    out_path = r"C:\Users\smcpa\Desktop\THE_IPO_PATH_Overview_v3.pptx"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    with open(out_path, "wb") as f:
        f.write(buf.read())
    print(f"SAVED: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
